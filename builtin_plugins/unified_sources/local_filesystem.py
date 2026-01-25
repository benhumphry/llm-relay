"""
Local Filesystem Unified Source Plugin.

A unified source for indexing local files mounted via Docker volumes,
with live document lookup capability.

Features:
- Document side: Index files from local directories for RAG
- Live side: Fetch full document content by filename
- Supports common document formats (text, markdown, PDF, Office, etc.)
- Configurable file filtering by extension
- Fuzzy filename matching for document resolution
"""

import logging
import os
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterator, Optional

from plugin_base.common import FieldDefinition, FieldType
from plugin_base.document_source import DocumentContent, DocumentInfo
from plugin_base.live_source import LiveDataResult, ParamDefinition
from plugin_base.unified_source import (
    MergeStrategy,
    PluginUnifiedSource,
    QueryAnalysis,
    QueryRouting,
)

logger = logging.getLogger(__name__)


class LocalFilesystemUnifiedSource(PluginUnifiedSource):
    """
    Local filesystem source with live document lookup.

    Indexes files from Docker-mounted directories for semantic search,
    and provides live lookup to fetch full document content by filename.
    """

    source_type = "local"
    display_name = "Local Filesystem"
    description = "Index files from local directories with full document lookup"
    category = "storage"
    icon = "📁"

    # Document store types this unified source handles
    handles_doc_source_types = ["local"]

    supports_rag = True
    supports_live = True  # Enable live document lookup
    supports_actions = False
    supports_incremental = True

    default_cache_ttl = 3600  # 1 hour - files don't change often

    _abstract = False

    @classmethod
    def get_designator_hint(cls) -> str:
        """Generate hint for designator prompt."""
        return (
            "Local filesystem with document lookup. Use action='lookup' with filename "
            "param when user asks for FULL CONTENT of a specific file (e.g., 'show me the invoice', "
            "'what's in document X'). Use action='list' to enumerate files. "
            "For semantic search across file contents, RAG will be used automatically."
        )

    @classmethod
    def build_config_from_store(cls, store) -> dict:
        """Build unified source config from a document store."""
        return {
            "path": store.source_path or "",
            "file_extensions": "",  # Use defaults
            "exclude_patterns": "",
        }

    # Supported file extensions for indexing
    DEFAULT_EXTENSIONS = {
        ".txt",
        ".md",
        ".markdown",
        ".rst",
        ".text",
        ".doc",
        ".docx",
        ".odt",
        ".xls",
        ".xlsx",
        ".ods",
        ".csv",
        ".ppt",
        ".pptx",
        ".odp",
        ".pdf",
        ".json",
        ".xml",
        ".yaml",
        ".yml",
        ".html",
        ".htm",
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".java",
        ".go",
        ".rs",
        ".cpp",
        ".c",
        ".h",
        ".rb",
        ".php",
        ".sh",
        ".bash",
        ".sql",
        ".graphql",
        ".log",
        ".ini",
        ".conf",
        ".cfg",
    }

    @classmethod
    def get_config_fields(cls) -> list[FieldDefinition]:
        """Configuration for admin UI."""
        return [
            FieldDefinition(
                name="path",
                label="Directory Path",
                field_type=FieldType.TEXT,
                required=True,
                help_text="Path to the directory to index (must be mounted in Docker)",
            ),
            FieldDefinition(
                name="recursive",
                label="Include Subdirectories",
                field_type=FieldType.BOOLEAN,
                default=True,
                help_text="Recursively index files in subdirectories",
            ),
            FieldDefinition(
                name="file_extensions",
                label="File Extensions",
                field_type=FieldType.TEXT,
                required=False,
                help_text="Comma-separated extensions to index (empty = all supported). Example: .txt,.md,.pdf",
            ),
            FieldDefinition(
                name="exclude_patterns",
                label="Exclude Patterns",
                field_type=FieldType.TEXT,
                required=False,
                help_text="Comma-separated patterns to exclude. Example: __pycache__,*.pyc,.git",
            ),
            FieldDefinition(
                name="max_file_size_mb",
                label="Max File Size (MB)",
                field_type=FieldType.INTEGER,
                default=10,
                help_text="Maximum file size to index in megabytes",
            ),
            FieldDefinition(
                name="index_schedule",
                label="Index Schedule",
                field_type=FieldType.SELECT,
                required=False,
                default="",
                options=[
                    {"value": "", "label": "Manual only"},
                    {"value": "0 */6 * * *", "label": "Every 6 hours"},
                    {"value": "0 0 * * *", "label": "Daily"},
                    {"value": "0 0 * * 0", "label": "Weekly"},
                ],
                help_text="How often to re-index files",
            ),
        ]

    @classmethod
    def get_live_params(cls) -> list[ParamDefinition]:
        """Parameters for live document lookup."""
        return [
            ParamDefinition(
                name="action",
                description="Query type: lookup (full document), list (enumerate files)",
                param_type="string",
                required=False,
                default="lookup",
                examples=["lookup", "list"],
            ),
            ParamDefinition(
                name="filename",
                description="Filename or partial name to lookup (for action=lookup)",
                param_type="string",
                required=False,
                examples=["invoice.pdf", "meeting notes", "report 2024"],
            ),
            ParamDefinition(
                name="max_results",
                description="Maximum files to list (for action=list)",
                param_type="integer",
                required=False,
                default=50,
            ),
        ]

    def __init__(self, config: dict):
        """Initialize with configuration."""
        self.path = config.get("path", "").rstrip("/")
        self.recursive = config.get("recursive", True)
        self.max_file_size_mb = config.get("max_file_size_mb", 10)
        self.index_schedule = config.get("index_schedule", "")

        # Parse file extensions
        ext_str = config.get("file_extensions", "")
        if ext_str:
            self.file_extensions = set(
                ext.strip().lower()
                if ext.strip().startswith(".")
                else f".{ext.strip().lower()}"
                for ext in ext_str.split(",")
                if ext.strip()
            )
        else:
            self.file_extensions = self.DEFAULT_EXTENSIONS

        # Parse exclude patterns
        exclude_str = config.get("exclude_patterns", "")
        self.exclude_patterns = (
            [p.strip() for p in exclude_str.split(",") if p.strip()]
            if exclude_str
            else []
        )

    def _should_exclude(self, path: str) -> bool:
        """Check if path matches any exclude pattern."""
        for pattern in self.exclude_patterns:
            if "*" in pattern:
                regex = pattern.replace(".", r"\.").replace("*", ".*")
                if re.search(regex, path):
                    return True
            else:
                if pattern in path:
                    return True
        return False

    def _get_mime_type(self, ext: str) -> str:
        """Get MIME type for file extension."""
        mime_map = {
            ".txt": "text/plain",
            ".md": "text/markdown",
            ".markdown": "text/markdown",
            ".rst": "text/x-rst",
            ".html": "text/html",
            ".htm": "text/html",
            ".xml": "application/xml",
            ".json": "application/json",
            ".yaml": "application/yaml",
            ".yml": "application/yaml",
            ".csv": "text/csv",
            ".pdf": "application/pdf",
            ".doc": "application/msword",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".xls": "application/vnd.ms-excel",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".ppt": "application/vnd.ms-powerpoint",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".py": "text/x-python",
            ".js": "text/javascript",
            ".ts": "text/typescript",
            ".java": "text/x-java-source",
            ".go": "text/x-go",
            ".rs": "text/x-rust",
            ".cpp": "text/x-c++src",
            ".c": "text/x-csrc",
            ".h": "text/x-chdr",
        }
        return mime_map.get(ext, "application/octet-stream")

    def list_documents(self) -> Iterator[DocumentInfo]:
        """Enumerate files for indexing."""
        if not os.path.isdir(self.path):
            logger.error(f"Directory not found: {self.path}")
            return

        logger.info(f"Listing local files in: {self.path}")
        total_files = 0
        max_size_bytes = self.max_file_size_mb * 1024 * 1024

        if self.recursive:
            for root, dirs, files in os.walk(self.path):
                dirs[:] = [
                    d for d in dirs if not self._should_exclude(os.path.join(root, d))
                ]
                for filename in files:
                    file_path = os.path.join(root, filename)
                    if self._should_exclude(file_path):
                        continue
                    doc = self._process_file(file_path, max_size_bytes)
                    if doc:
                        total_files += 1
                        yield doc
        else:
            for entry in os.scandir(self.path):
                if entry.is_file():
                    if self._should_exclude(entry.path):
                        continue
                    doc = self._process_file(entry.path, max_size_bytes)
                    if doc:
                        total_files += 1
                        yield doc

        logger.info(f"Found {total_files} files to index")

    def _process_file(
        self, file_path: str, max_size_bytes: int
    ) -> Optional[DocumentInfo]:
        """Process a single file for listing."""
        try:
            stat = os.stat(file_path)
            if stat.st_size > max_size_bytes:
                return None
            ext = Path(file_path).suffix.lower()
            if ext not in self.file_extensions:
                return None
            rel_path = os.path.relpath(file_path, self.path)
            modified_at = datetime.fromtimestamp(
                stat.st_mtime, tz=timezone.utc
            ).isoformat()
            mime_type = self._get_mime_type(ext)
            return DocumentInfo(
                uri=f"file://{file_path}",
                title=rel_path,
                mime_type=mime_type,
                modified_at=modified_at,
                metadata={
                    "path": file_path,
                    "relative_path": rel_path,
                    "size": stat.st_size,
                    "extension": ext,
                },
            )
        except Exception as e:
            logger.warning(f"Failed to process file {file_path}: {e}")
            return None

    def read_document(self, uri: str) -> Optional[DocumentContent]:
        """Read file content for indexing."""
        if not uri.startswith("file://"):
            return None
        file_path = uri.replace("file://", "")
        if not os.path.isfile(file_path):
            return None
        real_path = os.path.realpath(file_path)
        real_base = os.path.realpath(self.path)
        if not real_path.startswith(real_base):
            logger.error(f"File outside configured path: {file_path}")
            return None
        try:
            stat = os.stat(file_path)
            rel_path = os.path.relpath(file_path, self.path)
            ext = Path(file_path).suffix.lower()
            mime_type = self._get_mime_type(ext)
            if mime_type.startswith("text/") or mime_type in (
                "application/json",
                "application/xml",
                "application/yaml",
            ):
                with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
            else:
                with open(file_path, "rb") as f:
                    content = f.read()
            return DocumentContent(
                content=content,
                mime_type=mime_type,
                metadata={
                    "path": file_path,
                    "relative_path": rel_path,
                    "filename": os.path.basename(file_path),
                    "extension": ext,
                    "size": stat.st_size,
                    "source_type": "file",
                },
            )
        except Exception as e:
            logger.error(f"Failed to read file {file_path}: {e}")
            return None

    def fetch(self, params: dict) -> LiveDataResult:
        """Fetch full document content or list files."""
        action = params.get("action", "lookup")
        filename = params.get("filename", "")
        max_results = params.get("max_results", 50)
        two_pass_uri = params.get("_two_pass_uri")  # Direct URI from RAG metadata

        try:
            if action == "list":
                return self._fetch_list(max_results)
            elif action == "lookup" and (filename or two_pass_uri):
                return self._fetch_document(filename, two_pass_uri)
            else:
                return LiveDataResult(
                    success=False,
                    error="Specify action='lookup' with filename, or action='list'",
                )
        except Exception as e:
            logger.error(f"Filesystem fetch error: {e}")
            return LiveDataResult(success=False, error=str(e))

    def _fetch_list(self, max_results: int) -> LiveDataResult:
        """List files in the directory."""
        if not os.path.isdir(self.path):
            return LiveDataResult(
                success=False, error=f"Directory not found: {self.path}"
            )

        files = []
        for doc in self.list_documents():
            files.append(
                {
                    "name": doc.title,
                    "path": doc.metadata.get("path", ""),
                    "size": doc.metadata.get("size", 0),
                    "modified": doc.modified_at,
                }
            )
            if len(files) >= max_results:
                break

        formatted = self._format_file_list(files)
        return LiveDataResult(
            success=True,
            data=files,
            formatted=formatted,
            cache_ttl=self.default_cache_ttl,
        )

    def _fetch_document(
        self, filename: str, two_pass_uri: str = None
    ) -> LiveDataResult:
        """
        Fetch full document content by filename.

        Args:
            filename: Filename or partial name to search for
            two_pass_uri: Direct file path from RAG metadata (two-pass retrieval)
        """
        resolved_path = None

        # If two_pass_uri provided, use it directly (from RAG metadata)
        if two_pass_uri:
            # two_pass_uri may be the source_uri which could be file:// prefixed
            if two_pass_uri.startswith("file://"):
                direct_path = two_pass_uri[7:]  # Remove file:// prefix
            else:
                direct_path = two_pass_uri

            if os.path.isfile(direct_path):
                resolved_path = direct_path
                logger.debug(f"Two-pass: Direct path resolution: {resolved_path}")

        # Otherwise, use fuzzy matching
        if not resolved_path:
            # Build list of indexed files for fuzzy matching
            indexed_metadata = []
            for doc in self.list_documents():
                indexed_metadata.append(
                    {
                        "source_path": doc.metadata.get("path", ""),
                        "title": doc.title,
                        "name": os.path.basename(doc.metadata.get("path", "")),
                    }
                )

            # Use fuzzy matching to resolve filename
            resolved_path = self._resolve_document(filename, indexed_metadata)

        if not resolved_path:
            # Try direct path match as fallback
            direct_path = os.path.join(self.path, filename)
            if os.path.isfile(direct_path):
                resolved_path = direct_path
            else:
                return LiveDataResult(
                    success=False,
                    error=f"Could not find file matching '{filename}'",
                )

        # Read the document
        content_result = self.read_document(f"file://{resolved_path}")
        if not content_result:
            return LiveDataResult(
                success=False,
                error=f"Failed to read file: {resolved_path}",
            )

        # Extract text content
        text_content = self._extract_text_content(content_result, resolved_path)

        rel_path = os.path.relpath(resolved_path, self.path)
        formatted = f"### Full Content: {rel_path}\n\n{text_content}"

        return LiveDataResult(
            success=True,
            data={"path": resolved_path, "content": text_content},
            formatted=formatted,
            cache_ttl=self.default_cache_ttl,
        )

    def _resolve_document(
        self, query: str, indexed_metadata: list[dict]
    ) -> Optional[str]:
        """Resolve filename query to actual path using fuzzy matching."""
        try:
            from rapidfuzz import fuzz, process
        except ImportError:
            logger.warning("rapidfuzz not installed, using exact match")
            return self._exact_match_document(query, indexed_metadata)

        # Build candidates: {match_target: path}
        candidates: dict[str, str] = {}
        for meta in indexed_metadata:
            path = meta.get("source_path", "")
            if not path:
                continue

            # Use multiple match targets
            targets = [
                meta.get("title", ""),
                meta.get("name", ""),
                os.path.basename(path),
            ]
            for target in targets:
                if target:
                    target_lower = target.lower()
                    if target_lower not in candidates:
                        candidates[target_lower] = path

        if not candidates:
            return None

        # Fuzzy match - use partial_ratio for better substring matching
        query_lower = query.lower()
        match = process.extractOne(
            query_lower,
            candidates.keys(),
            scorer=fuzz.partial_ratio,
            score_cutoff=70,
        )

        if match:
            matched_target, score, _ = match
            path = candidates[matched_target]
            logger.info(f"Resolved '{query}' -> '{path}' (score={score:.0f})")
            return path

        return None

    def _exact_match_document(
        self, query: str, indexed_metadata: list[dict]
    ) -> Optional[str]:
        """Fallback exact substring matching."""
        query_lower = query.lower()
        for meta in indexed_metadata:
            path = meta.get("source_path", "")
            if not path:
                continue
            targets = [
                meta.get("title", ""),
                meta.get("name", ""),
                os.path.basename(path),
            ]
            for target in targets:
                if target and (
                    query_lower in target.lower() or target.lower() in query_lower
                ):
                    return path
        return None

    def _extract_text_content(
        self, content_result: DocumentContent, file_path: str
    ) -> str:
        """Extract text from document content, processing binary formats."""
        # If already text, return it
        if isinstance(content_result.content, str):
            return content_result.content

        # Binary content - need to process
        binary_content = content_result.content
        ext = Path(file_path).suffix.lower()
        mime_type = content_result.mime_type or ""

        # Try to extract text from binary formats
        try:
            # PDF
            if ext == ".pdf" or "pdf" in mime_type:
                return self._extract_pdf_text(binary_content)

            # DOCX
            if ext == ".docx" or "wordprocessingml" in mime_type:
                return self._extract_docx_text(binary_content)

            # XLSX
            if ext == ".xlsx" or "spreadsheetml" in mime_type:
                return self._extract_xlsx_text(binary_content)

            # PPTX
            if ext == ".pptx" or "presentationml" in mime_type:
                return self._extract_pptx_text(binary_content)

            # Unknown binary - return placeholder
            return f"[Binary file: {ext}, {len(binary_content)} bytes]"

        except Exception as e:
            logger.warning(f"Failed to extract text from {file_path}: {e}")
            return f"[Failed to extract text: {e}]"

    def _extract_pdf_text(self, binary_content: bytes) -> str:
        """Extract text from PDF."""
        try:
            import io

            import pypdf

            reader = pypdf.PdfReader(io.BytesIO(binary_content))
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            return "\n\n".join(text_parts) or "[No extractable text in PDF]"
        except ImportError:
            return "[pypdf not installed - cannot extract PDF text]"

    def _extract_docx_text(self, binary_content: bytes) -> str:
        """Extract text from DOCX including paragraphs and tables."""
        try:
            import io

            from docx import Document

            doc = Document(io.BytesIO(binary_content))
            parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)

            # Extract tables (many documents like invoices use tables)
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    # Skip empty rows
                    if any(cells):
                        table_rows.append(" | ".join(cells))
                if table_rows:
                    parts.append("\n".join(table_rows))

            return "\n\n".join(parts) if parts else "[No extractable text in DOCX]"
        except ImportError:
            return "[python-docx not installed - cannot extract DOCX text]"

    def _extract_xlsx_text(self, binary_content: bytes) -> str:
        """Extract text from XLSX."""
        try:
            import io

            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(binary_content), read_only=True)
            lines = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                lines.append(f"## Sheet: {sheet}")
                for row in ws.iter_rows(max_row=100, values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        lines.append(" | ".join(cells))
            return "\n".join(lines) or "[Empty spreadsheet]"
        except ImportError:
            return "[openpyxl not installed - cannot extract XLSX text]"

    def _extract_pptx_text(self, binary_content: bytes) -> str:
        """Extract text from PPTX."""
        try:
            import io

            from pptx import Presentation

            prs = Presentation(io.BytesIO(binary_content))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines.append(shape.text)
            return "\n".join(lines) or "[Empty presentation]"
        except ImportError:
            return "[python-pptx not installed - cannot extract PPTX text]"

    def _format_file_list(self, files: list[dict]) -> str:
        """Format file list for LLM context."""
        if not files:
            return "### Files\nNo files found in this directory."

        lines = ["### Files", f"Found {len(files)} file(s):\n"]
        for f in files:
            size_str = self._format_size(f.get("size", 0))
            lines.append(f"• **{f['name']}** ({size_str})")
            if f.get("modified"):
                lines.append(f"  Modified: {f['modified']}")
        return "\n".join(lines)

    def _format_size(self, size: int) -> str:
        """Format file size for display."""
        if size > 1024 * 1024:
            return f"{size / (1024 * 1024):.1f} MB"
        elif size > 1024:
            return f"{size / 1024:.1f} KB"
        return f"{size} bytes"

    def analyze_query(self, query: str, params: dict) -> QueryAnalysis:
        """Analyze query to determine RAG vs live lookup routing."""
        query_lower = query.lower()
        action = params.get("action", "")
        filename = params.get("filename", "")

        # Explicit lookup action -> Live only
        if action == "lookup" and filename:
            return QueryAnalysis(
                routing=QueryRouting.LIVE_ONLY,
                live_params=params,
                reason="Document lookup requested",
                max_live_results=1,
            )

        # Explicit list action -> Live only
        if action == "list":
            return QueryAnalysis(
                routing=QueryRouting.LIVE_ONLY,
                live_params=params,
                reason="File listing requested",
                max_live_results=params.get("max_results", 50),
            )

        # Detect document lookup intent from query
        # Two-pass: Use RAG to find relevant documents, then fetch full content
        lookup_patterns = [
            r"(full |entire |complete |whole )?content of",
            r"what('s| is| does) .* (say|contain)",
            r"show me .* (file|document)",
            r"read( me)? .* (file|document)",
            r"open .* (file|document)",
            r"(can you )?(get|fetch|retrieve|pull) .* (file|document)",
        ]

        import re

        for pattern in lookup_patterns:
            if re.search(pattern, query_lower):
                # Use TWO_PASS: RAG finds documents, then live fetches full content
                return QueryAnalysis(
                    routing=QueryRouting.TWO_PASS,
                    rag_query=query,
                    live_params={"action": "lookup", "filename": query},
                    merge_strategy=MergeStrategy.LIVE_FIRST,
                    reason="Document lookup intent - using two-pass (RAG finds, live fetches)",
                    max_rag_results=5,  # Find top 5 relevant documents
                    max_live_results=3,  # Fetch full content of top 3
                    two_pass_fetch_full=True,  # Enable full document fetch in pass 2
                )

        # Default -> RAG only for semantic search
        return QueryAnalysis(
            routing=QueryRouting.RAG_ONLY,
            rag_query=query,
            reason="Semantic search - using RAG index",
            max_rag_results=20,
        )

    def is_available(self) -> bool:
        """Check if directory is accessible."""
        return os.path.isdir(self.path)

    def test_connection(self) -> tuple[bool, str]:
        """Test filesystem access."""
        results = []
        if not os.path.isdir(self.path):
            return False, f"Directory not found: {self.path}"
        results.append(f"Directory: {self.path}")
        file_count = 0
        total_size = 0
        for doc in self.list_documents():
            file_count += 1
            total_size += doc.metadata.get("size", 0)
            if file_count >= 100:
                break
        results.append(
            f"Files: Found {file_count}{'+ ' if file_count >= 100 else ' '}file(s) to index"
        )
        if total_size > 1024 * 1024:
            size_str = f"{total_size / (1024 * 1024):.1f} MB"
        elif total_size > 1024:
            size_str = f"{total_size / 1024:.1f} KB"
        else:
            size_str = f"{total_size} bytes"
        results.append(f"Total size: {size_str}")
        return True, "\n".join(results)
